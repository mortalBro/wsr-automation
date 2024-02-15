# import mysql.connector
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN

# from utils import *

# connection = mysql.connector.connect(
#     host='localhost',
#     user='bhaiji',
#     password='triazine@123',
#     database='mortal'
# )


# cursor = connection.cursor()

# query = '''SELECT id, password, last_login, is_superuser, first_name, last_name, is_staff, is_active, date_joined, username, email, created_by, created_on, last_modified_on, last_modified_by, lock_unlock, status, user_type, first_login, profile_pic, locationcode, user_id, first_attempt, invalid_login, is_logedin, last_attempt, token, last_acive_time, raw_password, is_reset_password, last_reset_time, reset_password_date, report_view, sales_requisition, state_code, state_name, wd_ministartion
# FROM mortal.master_user
# WHERE id=5;'''

# cursor.execute(query)

# results = cursor.fetchall()

# presentation = Presentation('/home/tspl/Documents/wsrAutomation/WSR (Weekly Status Report).pptx')

# # Access the first slide (you may need to modify this based on your PowerPoint structure)
# # slide = presentation.slides[0]
# current_monday=current_monday()
# previous_monday=previous_monday()
# print(current_monday,previous_monday,"PPPPPPPPPPPPPPPPPPPPp")

# current_sunday=current_sunday()
# previous_sunday=previous_sunday()

# firstPageChange(presentation,previous_monday,current_monday)
# firstPageChange(presentation,previous_sunday,current_sunday)



# # left_inch = top_inch = Inches(1)  # Adjust the position as needed
# # width_inch = height_inch = Inches(2)  # Adjust the size as needed

# # textbox = slide.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
# # text_frame = textbox.text_frame
# # # for paragraph in text_frame.paragraphs:
# # #     print(paragraph.text)
# # text_frame.text = "Your text goes here"

# # # Access a specific shape on the slide (e.g., a text box)
# # textbox = slide.shapes[0].text_frame
# # for paragraph in textbox.paragraphs:
# #     print(paragraph.text)

# # print(textbox)

# # Update the text in the shape with data from the SQL query
# # for row in results:
# #     textbox.text += f"Column1: {row[0]}, Column2: {row[1]}\n"

# # # Save the modified presentation
# # presentation.save('updated_presentation.pptx')
# presentation.save('updated_presentation.pptx')
# connection.close()




# from datetime import datetime, timedelta

# def current_monday():
#     today = datetime.now()
#     days_to_subtract = today.weekday() + 7  
#     last_monday = today - timedelta(days=days_to_subtract)
#     return last_monday.strftime("%d/%m/%Y")

# def previous_monday():
#     today = datetime.now()
#     days_to_subtract = today.weekday() + 7 +7
#     last_monday = today - timedelta(days=days_to_subtract)
#     return last_monday.strftime("%d/%m/%Y")


# def current_sunday():
#     today = datetime.now()
#     days_to_subtract = today.weekday() + 7 - 6
#     last_sunday = today - timedelta(days=days_to_subtract)
#     return last_sunday.strftime("%d/%m/%Y")

# def previous_sunday():
#     today = datetime.now()
#     days_to_subtract = today.weekday() + 7 - 6
#     last_sunday = today - timedelta(days=days_to_subtract)
#     return last_sunday.strftime("%d/%m/%Y")


# #work in first slide
# def first_slide_work(presentation,current_monday,current_sunday,previous_monday,previous_sunday):
#   first_slide = presentation.slides[0]

#   # current_monday=current_monday()
#   # current_sunday=current_sunday()
#   # previous_monday=previous_monday()
#   # previous_sunday=previous_sunday()

#   for shape in first_slide.shapes:
#     if shape.has_text_frame and previous_monday in shape.text_frame.text:
#         existing_font = shape.text_frame.paragraphs[0].runs[0].font

#         new_paragraph = shape.text_frame.add_paragraph()
#         new_paragraph.text = current_monday

#         # Apply the font style to the new text
#         new_run = new_paragraph.runs[0]
#         new_run.font.size = existing_font.size
#         new_run.font.name = existing_font.name
#         new_run.font.bold = existing_font.bold
#         new_run.font.italic = existing_font.italic
#         new_run.font.color.theme_color = existing_font.color.theme_color
#         new_run.font.color.brightness = existing_font.color.brightness

#   for shape in first_slide.shapes:
#     if shape.has_text_frame and previous_sunday in shape.text_frame.text:
#         existing_font = shape.text_frame.paragraphs[0].runs[0].font

#         new_paragraph = shape.text_frame.add_paragraph()
#         new_paragraph.text = current_sunday

#         # Apply the font style to the new text
#         new_run = new_paragraph.runs[0]
#         new_run.font.size = existing_font.size
#         new_run.font.name = existing_font.name
#         new_run.font.bold = existing_font.bold
#         new_run.font.italic = existing_font.italic
#         new_run.font.color.theme_color = existing_font.color.theme_color
#         new_run.font.color.brightness = existing_font.color.brightness

from datetime import datetime, timedelta
import mysql.connector
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from countticket import *
from database import activeUserMirrorCalculation
from insertionInstance import insertionOfInstanceCount
import pandas as pd




presentation = Presentation('/home/tspl/Documents/wsrAutomation/WSR (Weekly Status Report).pptx')

first_slide = presentation.slides[2]
# Define the target string to identify the row
target_string = "Tickets Tracker(Mirror)\n(Resolved/ Logged)"
target_string_mt = "Tickets Tracker(MT)\n(Resolved/ Logged)"
active_mirror_user="Active Mirror Users"
active_mt_user="Active MT Users"



# Iterate through slides and shapes
for shape in first_slide.shapes:
  if shape.has_table:
    table = shape.table
    for row in table.rows:
      for cell in row.cells:
        if target_string in cell.text:
          val=countTicketReport().get('total')
          row.cells[1].text ="                                  "+f"{val}/{val}"

        elif target_string_mt in cell.text:
          val=countTicketReportMt().get('total')
          row.cells[1].text ="                                  "+f"{val}/{val}"

        elif active_mirror_user in cell.text:
          row.cells[1].text ="                                  "+str(12111)#activeUserMirrorCalculation()

        elif active_mt_user in cell.text:
          row.cells[1].text ="                                  "+str(112)  #


######################  &&                                             ----------------------------


# presentation = Presentation('/home/tspl/Documents/wsrAutomation/WSR (Weekly Status Report).pptx')
# fift_slide = presentation.slides[4]

# table_data = []
# for shape in fift_slide.shapes:
#     if shape.has_table:
#         table = shape.table
#         for row in table.rows:
#             row_data = [cell.text_frame.text for cell in row.cells]
#             table_data.append(row_data)

# headers = table_data[0]
# table_data = table_data[1:]

# # Create a pandas DataFrame
# df = pd.DataFrame(table_data, columns=headers)



# Print the DataFrame
# print(df)
################################################################                          ----------------------
# import pandas as pd
# from pptx import Presentation

# # Load the presentation
# presentation = Presentation('/home/tspl/Documents/wsrAutomation/WSR (Weekly Status Report).pptx')
# fift_slide = presentation.slides[4]

# # Extract table data from the PowerPoint slide
# table_data = []
# for shape in fift_slide.shapes:
#     if shape.has_table:
#         table = shape.table
#         for row in table.rows:
#             row_data = [cell.text_frame.text for cell in row.cells]
#             table_data.append(row_data)

# # Assuming you have headers, you can use the first row as column names
# headers = table_data[0]
# table_data = table_data[1:]

# # Create a pandas DataFrame
# df = pd.DataFrame(table_data, columns=headers)
# df['Mirror']="ndwhdgwd"

# # Make modifications to the DataFrame (example: changing a value)
# # df.loc[0, 'Column1'] = 'Modified Value'

# # Locate the shape containing the table
# table_shape = None
# for shape in fift_slide.shapes:
#     if shape.has_table:
#         table_shape = shape
#         break

# # Update the PowerPoint table with the modified DataFrame
# if table_shape:
#     table = table_shape.table
#     for row_idx, (_, row) in enumerate(df.iterrows()):
#         for col_idx, cell_value in enumerate(row):
#             table.cell(row_idx + 1, col_idx).text_frame.text = str(cell_value)

#     # Save the modified presentation
#     presentation.save('/home/tspl/Documents/wsrAutomation/Modified_WSR.pptx')
# else:
#     print("No table found in the specified slide.")


#-------------------------------------------------
import pandas as pd
from pptx import Presentation

# Load the presentation
presentation = Presentation('/home/tspl/Documents/wsrAutomation/WSR (Weekly Status Report).pptx')
fift_slide = presentation.slides[4]
table_data = []
for shape in fift_slide.shapes:
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            row_data = [cell.text_frame.text for cell in row.cells]
            table_data.append(row_data)

headers = table_data[0]
table_data = table_data[1:]

df = pd.DataFrame(table_data, columns=headers)



# Assuming "column_name" is the name of the column to which you want to add the prefix

df['Priority'] = df['Priority'].apply(lambda x: f"<b>{x}</b>")
df['Priority'] = df['Priority'].apply(lambda x: str(x).replace('<b>', '').replace('</b>', ''))
df['Priority'] = df['Priority'].apply(lambda x: f"       {x}")

df['CodNext'] = df['CodNext'].apply(lambda x: f"       {x}")
df['Esaarthi'] = df['Esaarthi'].apply(lambda x: f"       {x}")
df['GPI Website'] = df['GPI Website'].apply(lambda x: f"       {x}")
df['Mirror'] = df['Mirror'].apply(lambda x: f"       {x}")
df['MT'] = df['MT'].apply(lambda x: f"       {x}")
df['Total'] = df['Total'].apply(lambda x: f"       {x}")
df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: f"       {x}")

# df.loc[4] = df.loc[4].apply(lambda x: f"<b>{x}</b>")
# df.loc[4] = df.loc[4].apply(lambda x: str(x).replace('<b>', '').replace('</b>', ''))
# df.loc[4] = df.loc[4].apply(lambda x: f"       {x}")


df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: f"<b>{x}</b>")
df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: str(x).replace('<b>', '').replace('</b>', ''))
df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: f"       {x}")



# Locate the shape containing the table
table_shape = None
for shape in fift_slide.shapes:
    if shape.has_table:
        table_shape = shape
        break

# Update the PowerPoint table with the modified DataFrame
if table_shape:
    table = table_shape.table
    for row_idx, (_, row) in enumerate(df.iterrows()):
        for col_idx, cell_value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text_frame.text = str(cell_value)
            
            # Assuming "priority" is the column name you want to make bold
            if cell.text_frame.text == df.at[row_idx, 'Priority']:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                  
            if cell.text_frame.text == df.at[row_idx, 'Resolved Within SLA']:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
            


    # Save the modified presentation
    presentation.save('/home/tspl/Documents/wsrAutomation/Modified_WSR.pptx')
else:
    print("No table found in the specified slide.")


# for shape in fift_slide.shapes:
#   if shape.has_table:
#     table = shape.table
#     for row in table.rows:
#       for cell in row.cells:
#         print(cell.text)

        # if target_string in cell.text:
        #   row.cells[1].text ="gg"



# presentation.save("moo2n.pptx")





# for shape in first_slide.shapes:
#         # Check if the shape is a table
#   if shape.has_table:
#     table = shape.table
#     # Access data in the table
#     for row in table.rows:
#       for cell in row.cells:
#         print(cell.text)


