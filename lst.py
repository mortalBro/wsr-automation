import pandas as pd
from pptx import Presentation
from datetime import datetime, timedelta

def current_monday(k):
    today = datetime.now()
    days_to_subtract = today.weekday() + 7  +k
    last_monday = today - timedelta(days=days_to_subtract)
    return last_monday.strftime("%d-%m-%Y")


presentation = Presentation('/home/tspl/Documents/wsrAutomation/WSR (Weekly Status Report).pptx')
fift_slide = presentation.slides[12]
table_data = []
for shape in fift_slide.shapes:
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            row_data = [cell.text_frame.text for cell in row.cells]
            table_data.append(row_data)
headers = table_data[0]
table_data = table_data[1:]

df = pd.DataFrame(table_data, columns=headers)
print(df.iloc[0])

df.iloc[1][0]=current_monday(0)
df.iloc[2][0]=current_monday(1)
df.iloc[3][0]=current_monday(2)
df.iloc[4][0]=current_monday(3)
df.iloc[5][0]=current_monday(4)
df.iloc[6][0]=current_monday(5)
df.iloc[7][0]=current_monday(6)
# @@@@@/
first_row = df.iloc[2][1]
df.iloc[1][1]="MT LOGIN"
df.iloc[2][1]="MT LOGIN"
df.iloc[3][1]="MT LOGIN"
df.iloc[4][1]="MT LOGIN"
df.iloc[5][1]="MT LOGIN"
df.iloc[6][1]="MT LOGIN"
df.iloc[7][1]="MT LOGIN"



df.iloc[1][2]="transaction_MT"
df.iloc[2][2]="transaction_MT"
df.iloc[3][2]="transaction_MT"
df.iloc[4][2]="transaction_MT"
df.iloc[5][2]="transaction_MT"
df.iloc[6][2]="transaction_MT"
df.iloc[7][2]="transaction_MT"


df.iloc[1][3]="ss_login"
df.iloc[2][3]="ss_login"
df.iloc[3][3]="ss_login"
df.iloc[4][3]="ss_login"
df.iloc[5][3]="ss_login"
df.iloc[6][3]="ss_login"
df.iloc[7][3]="ss_login"



# Assuming "column_name" is the name of the column to which you want to add the prefix

# df['Priority'] = df['Priority'].apply(lambda x: f"<b>{x}</b>")
# df['Priority'] = df['Priority'].apply(lambda x: str(x).replace('<b>', '').replace('</b>', ''))
# df['Priority'] = df['Priority'].apply(lambda x: f"       {x}")

# df['CodNext'] = df['CodNext'].apply(lambda x: f"       {x}")
# df['Esaarthi'] = df['Esaarthi'].apply(lambda x: f"       {x}")
# df['GPI Website'] = df['GPI Website'].apply(lambda x: f"       {x}")
# df['Mirror'] = df['Mirror'].apply(lambda x: f"       {x}")
# df['MT'] = df['MT'].apply(lambda x: f"       {x}")
# df['Total'] = df['Total'].apply(lambda x: f"       {x}")
# df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: f"       {x}")

# # df.loc[4] = df.loc[4].apply(lambda x: f"<b>{x}</b>")
# # df.loc[4] = df.loc[4].apply(lambda x: str(x).replace('<b>', '').replace('</b>', ''))
# # df.loc[4] = df.loc[4].apply(lambda x: f"       {x}")


# df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: f"<b>{x}</b>")
# df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: str(x).replace('<b>', '').replace('</b>', ''))
# df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: f"       {x}")

table_shape = None
for shape in fift_slide.shapes:
    if shape.has_table:
        table_shape = shape
        break

# Update the PowerPoint table with the modified DataFrame
if table_shape:
    table = table_shape.table
    for row_idx, (_, row) in enumerate(df.iterrows()):
        for col_idx, cell_value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text_frame.text = str(cell_value)
            
            if row_idx == 0:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True


    # Save the modified presentation
    presentation.save('/home/tspl/Documents/wsrAutomation/Modified_WSR.pptx')
else:
    print("No table found in the specified slide.")

