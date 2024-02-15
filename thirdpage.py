from datetime import datetime, timedelta
import mysql.connector
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from countticket import *
from database import activeUserMirrorCalculation
from insertionInstance import insertionOfInstanceCount
import pandas as pd



def thirdPageWork():
  presentation = Presentation('/home/tspl/Documents/wsrAutomation/WSR (Weekly Status Report).pptx')

  first_slide = presentation.slides[2]
  # Define the target string to identify the row
  target_string = "Tickets Tracker(Mirror)\n(Resolved/ Logged)"
  target_string_mt = "Tickets Tracker(MT)\n(Resolved/ Logged)"
  active_mirror_user="Active Mirror Users"
  active_mt_user="Active MT Users"

  # Iterate through slides and shapes
  for shape in first_slide.shapes:
    if shape.has_table:
      table = shape.table
      for row in table.rows:
        for cell in row.cells:
          if target_string in cell.text:
            val=countTicketReport().get('total')
            row.cells[1].text ="                                  "+f"{val}/{val}"

          elif target_string_mt in cell.text:
            val=countTicketReportMt().get('total')
            row.cells[1].text ="                                  "+f"{val}/{val}"

          elif active_mirror_user in cell.text:
            row.cells[1].text ="                                  "+str(12111)#activeUserMirrorCalculation()

          elif active_mt_user in cell.text:
            row.cells[1].text ="                                  "+str(112)  #