import pandas as pd
from pptx import Presentation

def fourthPage():
  presentation = Presentation('/home/tspl/Documents/wsrAutomation/WSR (Weekly Status Report).pptx')
  fift_slide = presentation.slides[4]
  table_data = []
  for shape in fift_slide.shapes:
      if shape.has_table:
          table = shape.table
          for row in table.rows:
              row_data = [cell.text_frame.text for cell in row.cells]
              table_data.append(row_data)

  headers = table_data[0]
  table_data = table_data[1:]

  df = pd.DataFrame(table_data, columns=headers)
  # Assuming "column_name" is the name of the column to which you want to add the prefix

  df['Priority'] = df['Priority'].apply(lambda x: f"<b>{x}</b>")
  df['Priority'] = df['Priority'].apply(lambda x: str(x).replace('<b>', '').replace('</b>', ''))
  df['Priority'] = df['Priority'].apply(lambda x: f"       {x}")

  df['CodNext'] = df['CodNext'].apply(lambda x: f"       {x}")
  df['Esaarthi'] = df['Esaarthi'].apply(lambda x: f"       {x}")
  df['GPI Website'] = df['GPI Website'].apply(lambda x: f"       {x}")
  df['Mirror'] = df['Mirror'].apply(lambda x: f"       {x}")
  df['MT'] = df['MT'].apply(lambda x: f"       {x}")
  df['Total'] = df['Total'].apply(lambda x: f"       {x}")
  df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: f"       {x}")

  # df.loc[4] = df.loc[4].apply(lambda x: f"<b>{x}</b>")
  # df.loc[4] = df.loc[4].apply(lambda x: str(x).replace('<b>', '').replace('</b>', ''))
  # df.loc[4] = df.loc[4].apply(lambda x: f"       {x}")


  df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: f"<b>{x}</b>")
  df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: str(x).replace('<b>', '').replace('</b>', ''))
  df['Resolved Within SLA'] = df['Resolved Within SLA'].apply(lambda x: f"       {x}")

  # Locate the shape containing the table
  table_shape = None
  for shape in fift_slide.shapes:
      if shape.has_table:
          table_shape = shape
          break

  # Update the PowerPoint table with the modified DataFrame
  if table_shape:
      table = table_shape.table
      for row_idx, (_, row) in enumerate(df.iterrows()):
          for col_idx, cell_value in enumerate(row):
              cell = table.cell(row_idx + 1, col_idx)
              cell.text_frame.text = str(cell_value)
              # Assuming "priority" is the column name you want to make bold
              if cell.text_frame.text == df.at[row_idx, 'Priority']:
                  for paragraph in cell.text_frame.paragraphs:
                      for run in paragraph.runs:
                          run.font.bold = True
                    
              if cell.text_frame.text == df.at[row_idx, 'Resolved Within SLA']:
                  for paragraph in cell.text_frame.paragraphs:
                      for run in paragraph.runs:
                          run.font.bold = True
              


      # Save the modified presentation
      presentation.save('/home/tspl/Documents/wsrAutomation/Modified_WSR.pptx')
  else:
      print("No table found in the specified slide.")